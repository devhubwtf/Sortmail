"""
Attachment Extractor
--------------------
Extracts and stores attachments from emails.
Async-safe: uses aiofiles for all I/O.
"""

import os
import aiofiles
from typing import List
from datetime import datetime, timezone

from app.config import settings
from contracts import AttachmentRef


async def extract_attachments(
    thread_id: str,
    raw_attachments: List[dict],
) -> List[AttachmentRef]:
    """
    Extract attachments from email data and store them.
    
    Args:
        thread_id: Parent thread ID
        raw_attachments: Raw attachment data from email API
        
    Returns:
        List of AttachmentRef with storage paths
    """
    results = []
    
    for i, att in enumerate(raw_attachments):
        attachment_id = f"att-{thread_id}-{i}"
        original_filename = att.get("filename", f"attachment_{i}")
        
        # Generate smart filename
        smart_filename = _generate_smart_filename(
            original_filename,
            att.get("content_type", "application/octet-stream"),
        )
        
        # Store file — expects (user_id, filename, data)
        # Here we use attachment_id as user_id placeholder since we don't
        # have user_id in this legacy interface. Use sync_service flow for
        # proper user-scoped storage.
        storage_path = await _store_attachment(
            attachment_id,
            smart_filename,
            att.get("data", b""),
        )
        
        results.append(AttachmentRef(
            attachment_id=attachment_id,
            message_id=att.get("message_id", ""),
            filename=smart_filename,
            original_filename=original_filename,
            mime_type=att.get("content_type", "application/octet-stream"),
            storage_path=storage_path,
            size_bytes=len(att.get("data", b"")),
        ))
    
    return results


def _generate_smart_filename(original: str, mime_type: str) -> str:
    """
    Generate a contextual filename.
    
    Example: 'scan001.pdf' -> 'Document_2026-01-18.pdf'
    """
    # TODO: Implement smart renaming based on content analysis
    # For now, keep original with date prefix
    date_prefix = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    name, ext = os.path.splitext(original)
    return f"{name}_{date_prefix}{ext}"


async def _store_attachment(
    user_id: str,
    filename: str,
    data: bytes,
) -> str:
    """
    Store attachment to configured storage directory.

    Args:
        user_id: User ID — used to create a per-user subdirectory for isolation.
        filename: Filename to save the file as (should already be sanitized).
        data: Raw bytes to write.

    Returns:
        Absolute path where file was stored.
    """
    # Per-user directory for isolation
    storage_dir = os.path.join(settings.STORAGE_PATH, str(user_id))
    import asyncio
    await asyncio.to_thread(os.makedirs, storage_dir, exist_ok=True)
    
    # Sanitize filename (strip path separators to prevent directory traversal)
    safe_filename = os.path.basename(filename) or "unnamed_attachment"
    file_path = os.path.join(storage_dir, safe_filename)
    
    # Async write
    async with aiofiles.open(file_path, "wb") as f:
        await f.write(data)
    
    return file_path


SUPPORTED_MIME_TYPES = [
    "application/pdf",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
]


def is_supported_attachment(mime_type: str) -> bool:
    """Check if attachment type is supported for processing."""
    if not mime_type: return False
    return mime_type.lower().strip() in SUPPORTED_MIME_TYPES

async def index_attachment(attachment_id: str, user_id: str, db):
    """
    Extract text from a downloaded attachment and route it via the RAG Hybrid Context Strategy.
    """
    from sqlalchemy.ext.asyncio import AsyncSession
    from sqlalchemy import select
    from models.attachment import Attachment
    import logging
    
    logger = logging.getLogger(__name__)

    stmt = select(Attachment).where(Attachment.id == attachment_id)
    result = await db.execute(stmt)
    att = result.scalars().first()
    
    if not att or not att.storage_path:
        logger.warning(f"Attachment {attachment_id} missing or lacks storage path.")
        return
        
    logger.info(f"Extracting text from attachment: {att.filename}")
    text = ""
    
    try:
        if "pdf" in att.mime_type.lower():
            text = await _extract_pdf(att.storage_path)
        elif "word" in att.mime_type.lower() or "document" in att.mime_type.lower():
            text = await _extract_docx(att.storage_path)
        elif "powerpoint" in att.mime_type.lower() or "presentation" in att.mime_type.lower():
            text = await _extract_pptx(att.storage_path)
    except Exception as e:
        logger.error(f"Text extraction failed for {att.filename}: {e}")
        
    if text:
        # Pre-allocate for the upcoming strategy call
        att.extracted_text = text
        await db.commit()
        
        # Hook into intelligent router strategy
        try:
            from core.rag.attachment_strategy import AttachmentContextStrategy
            await AttachmentContextStrategy.process_attachment(att, user_id, db)
        except Exception as e:
            logger.error(f"Context routing failed for {att.filename}: {e}")

async def _extract_pdf(path: str) -> str:
    import asyncio
    def _read():
        from pypdf import PdfReader
        reader = PdfReader(path)
        return "\n".join([page.extract_text() or "" for page in reader.pages])
    return await asyncio.to_thread(_read)

async def _extract_docx(path: str) -> str:
    import asyncio
    def _read():
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    return await asyncio.to_thread(_read)

async def _extract_pptx(path: str) -> str:
    import asyncio
    def _read():
        from pptx import Presentation
        prs = Presentation(path)
        return "\n".join([
            shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
        ])
    return await asyncio.to_thread(_read)
