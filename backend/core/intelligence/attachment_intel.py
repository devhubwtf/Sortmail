"""
Attachment Intelligence
-----------------------
Analyzes document attachments (PDF, DOCX, PPTX) and produces deep analysis.
Uses Gemini Flash to extract summary, key points, action items, financial data, etc.
"""

import json
import logging
from typing import Optional, Dict, Any
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.future import select

from models.attachment import Attachment, AttachmentStatus
import google.generativeai as genai
from app.config import settings

logger = logging.getLogger(__name__)

# Configure Gemini once globally
if getattr(settings, "GEMINI_API_KEY", None):
    genai.configure(api_key=settings.GEMINI_API_KEY)


async def analyze_attachment(attachment_id: str, db: AsyncSession) -> Optional[Dict[str, Any]]:
    """
    Extracts text from a downloaded attachment and runs a unified Gemini Flash analysis.
    Saves the result to attachment.intel_json and returns the JSON payload.
    """
    stmt = select(Attachment).where(Attachment.id == attachment_id)
    result = await db.execute(stmt)
    attachment = result.scalars().first()
    
    if not attachment:
        logger.warning(f"Attachment {attachment_id} not found.")
        return None
        
    if not attachment.storage_path:
        logger.warning(f"Attachment {attachment_id} has no storage path yet.")
        return None

    # 1. Extract Document Text
    text = await _extract_text(attachment)
    if not text or len(text.strip()) < 10:
        logger.debug(f"Attachment {attachment_id} text extraction failed or insufficient.")
        attachment.status = AttachmentStatus.FAILED
        await db.commit()
        return None
        
    # Chunk text broadly if it's massive to avoid absolutely destroying tokens
    # Gemini Flash supports 1m tokens, but good to cap at ~50k characters for speed
    safe_text = text[:50000]

    # 2. Build Deep Analysis Prompt
    prompt = f"""
You are an expert document analyst. Analyze this business document and return pure JSON with no markdown formatting.
Ensure the keys match exactly:
- `document_type`: one of [contract, invoice, proposal, report, resume, meeting_notes, other]
- `summary`: 2-3 sentence overview
- `key_points`: list of 3-5 important facts
- `action_items`: list of specific things the user must do
- `financial_amounts`: any monetary values mentioned, e.g., ["$1,500 total due", "€50k budget"]
- `dates_and_deadlines`: important dates mentioned
- `risk_flags`: anything unusual or requiring attention (empty list if none)
- `parties_involved`: list of companies or people named

DOCUMENT TEXT TO ANALYZE:
---
{safe_text}
---

Return ONLY valid JSON.
"""
    
    try:
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = await model.generate_content_async(
            contents=prompt,
            generation_config=genai.types.GenerationConfig(
                temperature=0.1
            )
        )
        
        raw_json = response.text
        intel_dict = json.loads(raw_json)
        
        # 3. Save back to DB
        attachment.intel_json = intel_dict
        attachment.extracted_text = safe_text[:2000] # Save a snippet
        attachment.status = AttachmentStatus.INDEXED
        
        await db.commit()
        logger.info(f"✅ Generated intelligence for attachment {attachment.id}")
        return intel_dict
        
    except Exception as e:
        logger.error(f"Failed to analyze attachment {attachment.id}: {e}")
        attachment.status = AttachmentStatus.FAILED
        await db.commit()
        return None


# ─────────────────────────────────────────────────────────────────────────────
# Text Extractors
# ─────────────────────────────────────────────────────────────────────────────

async def _extract_text(attachment: Attachment) -> str:
    """Route text extraction based on mime type or extension."""
    try:
        mime = (attachment.mime_type or "").lower()
        path = attachment.storage_path
        
        if mime == "application/pdf" or path.endswith(".pdf"):
            return _extract_pdf_text(path)
        elif "wordprocessing" in mime or path.endswith(".docx"):
            return _extract_docx_text(path)
        elif "presentation" in mime or path.endswith(".pptx"):
            return _extract_pptx_text(path)
        elif mime.startswith("text/"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        else:
            logger.debug(f"Unsupported mime type for text extraction: {mime}")
            return ""
    except Exception as e:
        logger.error(f"Text extraction failed for {attachment.id}: {e}")
        return ""


def _extract_pdf_text(path: str) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        return "\\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
    except ImportError:
        logger.warning("pypdf not installed.")
        return ""
    except Exception as e:
        logger.warning(f"Failed parsing PDF {path}: {e}")
        return ""


def _extract_docx_text(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        return "\\n".join([p.text for p in doc.paragraphs if p.text])
    except ImportError:
        logger.warning("python-docx not installed.")
        return ""
    except Exception as e:
        logger.warning(f"Failed parsing DOCX {path}: {e}")
        return ""


def _extract_pptx_text(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        return "\\n".join(text)
    except ImportError:
        logger.warning("python-pptx not installed.")
        return ""
    except Exception as e:
        logger.warning(f"Failed parsing PPTX {path}: {e}")
        return ""
